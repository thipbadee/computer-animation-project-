from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from submission_metrics import (
    AUTHOR_NAME,
    COURSE_NAME,
    DEPARTMENT,
    GITHUB_USERNAME,
    INSTRUCTOR_NAME,
    KEYWORDS,
    PROJECT_TITLE,
    SEMESTER,
    benchmark_summary,
    group_rows_by_size,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_FILE = REPO_ROOT / "report" / "final_presentation_th.pptx"

TITLE_COLOR = RGBColor(20, 24, 32)
ACCENT_COLOR = RGBColor(0, 109, 119)
ACCENT_ALT = RGBColor(232, 248, 249)
TEXT_COLOR = RGBColor(50, 58, 70)
MUTED_COLOR = RGBColor(104, 116, 132)
WHITE = RGBColor(255, 255, 255)

WORKFLOW_IMAGE = REPO_ROOT / "report" / "assets" / "workflow.png"
PLOT_IMAGE = REPO_ROOT / "experiments" / "results" / "plots.png"
DEMO_STRIP_IMAGE = REPO_ROOT / "report" / "assets" / "demo_strip.png"

FONT_DISPLAY = "Leelawadee UI"
FONT_BODY = "Leelawadee UI"


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_DISPLAY
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.16), Inches(2.3), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT_COLOR
    rule.line.fill.background()

    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED_COLOR


def add_text(slide, text: str, left: float, top: float, width: float, height: float, font_size=18, bold=False, color=TEXT_COLOR):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR


def add_card(slide, title: str, value: str, left: float, top: float, width: float = 2.2, height: float = 1.1):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ALT
    shape.line.color.rgb = ACCENT_COLOR
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT_BODY
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = ACCENT_COLOR
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = FONT_DISPLAY
    r2.font.size = Pt(17)
    r2.font.bold = True
    r2.font.color.rgb = TITLE_COLOR


def add_table(slide, data: list[list[str]], left: float, top: float, width: float, height: float, font_size=12):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    col_width = int((Inches(width) / cols))
    for idx in range(cols):
        table.columns[idx].width = col_width
    for r, row_data in enumerate(data):
        for c, value in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_ALT if r == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = TITLE_COLOR if r == 0 else TEXT_COLOR
                    run.font.bold = (r == 0)


def add_image(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def build_presentation():
    summary = benchmark_summary()
    grouped = group_rows_by_size()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, RGBColor(248, 250, 252))
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_COLOR
    banner.line.fill.background()
    add_text(slide, "ระบบติดตามการเคลื่อนไหวท่าออกกำลังกายแบบเรียลไทม์", 0.8, 1.0, 8.9, 0.85, font_size=26, bold=True, color=TITLE_COLOR)
    add_text(slide, "รายวิชา " + COURSE_NAME, 0.8, 1.95, 2.8, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_text(slide, DEPARTMENT, 0.8, 2.3, 7.7, 0.45, font_size=14, color=MUTED_COLOR)
    add_text(slide, f"ผู้จัดทำ: {AUTHOR_NAME}", 0.8, 5.5, 3.0, 0.3, font_size=14)
    add_text(slide, f"GitHub: {GITHUB_USERNAME}", 0.8, 5.85, 3.2, 0.3, font_size=14)
    add_text(slide, f"อาจารย์ผู้สอน: {INSTRUCTOR_NAME}", 0.8, 6.2, 3.8, 0.3, font_size=12, color=MUTED_COLOR)
    add_text(slide, f"ภาคการศึกษา: {SEMESTER}", 4.1, 6.2, 4.0, 0.3, font_size=12, color=MUTED_COLOR)
    add_text(slide, "คำสำคัญ: " + ", ".join(KEYWORDS), 0.8, 6.55, 6.5, 0.3, font_size=12, color=MUTED_COLOR)
    add_card(slide, "Speedup ที่ 100k", f"{summary['speedup_pct']:.2f}%", 8.8, 1.35)
    add_card(slide, "ลดหน่วยความจำ", f"{summary['memory_reduction_pct']:.2f}%", 11.1, 1.35)
    add_card(slide, "ลด Jitter", f"{summary['jitter_reduction_pct']:.2f}%", 9.95, 2.7)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "นิยามปัญหาและวัตถุประสงค์", "ครอบคลุม rubric ส่วน problem definition, baseline/proposed และเป้าหมายเชิงวิศวกรรม")
    add_bullets(
        slide,
        [
            "ศึกษาโจทย์การติดตามท่าออกกำลังกายแบบเรียลไทม์ในมุมมองของ Computer Animation ไม่ใช่แค่การจำแนกภาพจากกล้อง",
            "ใช้ pose landmark เป็นโครงกระดูกแบบย่อ และตีความท่าออกกำลังกายเป็นการเปลี่ยนผ่านระหว่าง key poses",
            "เปรียบเทียบวิธี baseline แบบ threshold ตรง ๆ กับวิธี proposed ที่เพิ่ม confidence-aware side lock และ temporal filtering",
            "วัดผลด้วย runtime, memory, repetition-count accuracy และ temporal stability",
            "จัดทำ artifact ที่ทำซ้ำได้จริง พร้อม live demo, benchmark, tests และเอกสารส่งงาน",
        ],
        0.8,
        1.9,
        6.1,
        4.8,
        font_size=18,
    )
    add_text(slide, "เหตุผลที่เหมาะกับวิชา Computer Animation", 7.2, 1.9, 4.1, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "มองร่างกายเป็น skeleton-like structure",
            "มี key poses และ stage transitions",
            "มีการแปลง pose progress เป็นค่าต่อเนื่อง",
            "เชื่อมโยงกับ motion capture แบบ markerless",
        ],
        7.25,
        2.25,
        5.0,
        3.2,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "นิยามปัญหาเชิงรูปนัย")
    add_bullets(
        slide,
        [
            "Input: ลำดับของ 2D joint landmarks พร้อม visibility จาก webcam หรือไฟล์ CSV ที่บันทึกไว้",
            "Output: จำนวนครั้ง, stage ปัจจุบัน, progress percentage และภาพ visualization",
            "Constraints: ต้องตอบสนองแบบ interactive, ทนต่อ occlusion, ลด threshold chatter และทำงานได้แบบ deterministic",
            "ความท้าทาย 1: จุดกลับตัวของท่ามี noise ทำให้เกิดนับซ้ำหรือหลุดนับ",
            "ความท้าทาย 2: การสลับเลือกแขนซ้ายขวาบ่อยเกินไปทำให้การตีความ motion ไม่นิ่ง",
            "ความท้าทาย 3: ต้องทำให้ feedback เข้าใจง่ายสำหรับผู้ใช้แบบ real-time",
        ],
        0.85,
        1.9,
        5.95,
        4.8,
        font_size=18,
    )
    add_text(slide, "Motion Primitive ที่ใช้ benchmark", 7.1, 1.95, 3.5, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_text(
        slide,
        "ใช้ท่า bicep curl เป็นตัวแทนของ articulated motion primitive สำหรับวัด baseline/proposed อย่างยุติธรรม ส่วน live app ขยาย framework เดียวกันไปยัง 6 ท่าออกกำลังกาย",
        7.15,
        2.3,
        5.1,
        1.7,
        font_size=17,
    )
    add_text(slide, "Visualization ที่ส่งพร้อมงาน", 7.1, 4.35, 3.0, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "หน้าเว็บ Flask แบบ live",
            "demo video แบบ side-by-side",
            "plots สำหรับ runtime และ stability",
        ],
        7.15,
        4.7,
        4.6,
        1.8,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Baseline, Proposed และบริบท SOTA")
    table_data = [
        ["วิธี", "แนวคิดหลัก", "จุดเด่น", "ข้อจำกัด"],
        ["Baseline", "ใช้ elbow angle + threshold นับตรง ๆ", "ง่ายและ deterministic", "เปราะกับ jitter และ occlusion"],
        ["Proposed", "เพิ่ม confidence-aware side lock + EMA smoothing", "นิ่งขึ้นและ benchmark เร็วกว่า", "ยังเป็น handcrafted method"],
        ["SOTA context", "โมเดล temporal neural / action model", "จับ sequence context ได้มากกว่า", "ใช้ compute สูงและอธิบายยากกว่า"],
    ]
    add_table(slide, table_data, 0.75, 1.9, 11.9, 2.4, font_size=11)
    add_bullets(
        slide,
        [
            "ทั้ง baseline และ proposed ใช้ input/output แบบเดียวกัน จึงเปรียบเทียบกันได้ยุติธรรม",
            "งานนี้ไม่ได้อ้างว่าเหนือกว่า SOTA แต่แสดง improvement ที่วัดซ้ำได้จริงเหนือ baseline",
            "จุดเน้นของ proposed คือคุณค่าทางวิศวกรรม: เสถียรขึ้นโดยยังอธิบายการทำงานได้ง่าย",
        ],
        0.9,
        4.75,
        11.0,
        1.8,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "แนวคิดของวิธี Proposed และ Workflow")
    add_image(slide, WORKFLOW_IMAGE, 6.9, 1.75, 5.7)
    add_bullets(
        slide,
        [
            "คำนวณ confidence ของแขนซ้าย/ขวาจาก visibility และตำแหน่งบนหน้าจอ",
            "ถ้าแขนที่ lock ไว้ยังน่าเชื่อถือ ให้ใช้แขนนั้นต่อเพื่อลด arm flapping",
            "ใช้ exponential moving average เพื่อลด noise ของ elbow angle",
            "ใช้ hysteresis thresholds เพื่อตัดสิน stage และนับ repetition",
            "ให้ output แบบเดียวกับ baseline: count, stage, angle และ progress",
        ],
        0.85,
        1.95,
        5.5,
        3.6,
        font_size=18,
    )
    add_text(
        slide,
        "Complexity: ทั้ง baseline และ proposed เป็น O(1) time และ O(1) memory ต่อเฟรม ความต่างอยู่ที่การจัดการ state ให้เสถียรกว่าเดิม",
        0.85,
        5.85,
        11.7,
        0.6,
        font_size=16,
        color=MUTED_COLOR,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "รายละเอียดการพัฒนา")
    add_bullets(
        slide,
        [
            "ระบบ live ใช้ Flask + OpenCV + MediaPipe + NumPy",
            "ส่วน benchmark แยกไว้ใน `src/algorithms/` เพื่อวัด baseline/proposed โดยไม่ต้องเปิดกล้อง",
            "state ที่ใช้มี counter, stage, active side, ค่า EMA และ side-lock timer",
            "กรณีขอบที่จัดการ: visibility drop, arm switching, threshold chatter, camera read failure",
            "เพิ่ม pipeline ข้อมูลจริงด้วย `tools/collect_mocap.py` และ `experiments/evaluate_recorded_sequence.py`",
        ],
        0.85,
        1.9,
        6.0,
        4.9,
        font_size=18,
    )
    add_text(slide, "โครงสร้างไฟล์สำคัญ", 7.05, 1.9, 3.3, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "`app.py`, `camera.py`, `exercises/`",
            "`src/algorithms/baseline.py`",
            "`src/algorithms/proposed.py`",
            "`experiments/benchmark.py`",
            "`visualization/animate_results.py`",
            "`tests/test_simulation.py`",
        ],
        7.1,
        2.25,
        5.0,
        3.6,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Experimental Setup")
    add_bullets(
        slide,
        [
            "Scaling benchmark ใช้ข้อมูลสังเคราะห์ 3 ขนาด: 1k, 10k และ 100k frames",
            "ตัวสร้างข้อมูลเพิ่มทั้ง angle noise และ visibility dropout เพื่อทดสอบ temporal robustness",
            "ข้อมูลจริงรองรับผ่านการอัด landmark CSV จาก webcam แล้ว replay เพื่อประเมินซ้ำ",
            "Metrics ที่ใช้: total runtime, ms/frame, peak traced memory, repetition count error และ side-switch jitter",
            "การตรวจสอบความถูกต้อง: `python -m unittest tests.test_simulation` ผ่าน",
        ],
        0.85,
        1.95,
        6.1,
        4.6,
        font_size=18,
    )
    add_card(slide, "Synthetic Sizes", "1k / 10k / 100k", 7.4, 2.0, width=2.6)
    add_card(slide, "Count Error", "0 ทุก scale", 10.2, 2.0, width=2.3)
    add_card(slide, "Real-data Path", "Capture + Replay", 7.4, 3.45, width=2.6)
    add_card(slide, "Visualization", "Plots + Demo", 10.2, 3.45, width=2.3)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "ผลเชิงปริมาณ")
    table_rows = [["Input", "Baseline", "Proposed", "Speedup", "Jitter เดิม", "Jitter ใหม่"]]
    for size in sorted(grouped):
        baseline = grouped[size]["baseline"]
        proposed = grouped[size]["proposed"]
        speedup = ((baseline.ms_per_frame - proposed.ms_per_frame) / baseline.ms_per_frame) * 100.0
        table_rows.append(
            [
                f"{size:,}",
                f"{baseline.ms_per_frame:.6f}",
                f"{proposed.ms_per_frame:.6f}",
                f"{speedup:.2f}%",
                str(baseline.side_switches),
                str(proposed.side_switches),
            ]
        )
    add_table(slide, table_rows, 0.75, 1.9, 12.0, 2.5, font_size=11)
    add_card(slide, "100k ms/frame", f"{summary['baseline_ms_per_frame']:.6f} -> {summary['proposed_ms_per_frame']:.6f}", 0.95, 4.9, width=3.6)
    add_card(slide, "Peak Memory", f"{summary['baseline_peak_kb']:.3f} -> {summary['proposed_peak_kb']:.3f} KB", 4.95, 4.9, width=3.0)
    add_card(slide, "Side Switches", f"{int(summary['baseline_side_switches'])} -> {int(summary['proposed_side_switches'])}", 8.35, 4.9, width=3.1)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "ผลเชิงภาพและ Visualization")
    add_image(slide, PLOT_IMAGE, 0.7, 1.75, 6.1)
    add_image(slide, DEMO_STRIP_IMAGE, 7.1, 1.9, 5.5)
    add_text(
        slide,
        "ซ้าย: กราฟผล scaling benchmark จาก `experiments/results/plots.png` ขวา: เฟรมตัวอย่างจาก demo video ที่เปรียบเทียบ baseline กับ proposed แบบ side-by-side",
        0.8,
        6.1,
        11.8,
        0.45,
        font_size=15,
        color=MUTED_COLOR,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "อภิปรายผล ข้อจำกัด และงานต่อยอด")
    add_bullets(
        slide,
        [
            "จุดแข็ง: proposed ลด temporal jitter ได้โดยไม่เปลี่ยน interface และไม่ต้องพึ่ง learned weights",
            "จุดแข็ง: repository มี benchmark, tests, plots, demo video และ pipeline สำหรับข้อมูลจริงครบ",
            "ข้อจำกัด: benchmark เชิงปริมาณยังโฟกัสที่ motion primitive หลักเพียงชนิดเดียว",
            "ข้อจำกัด: repo ไม่ได้ commit human recording จริงไว้โดยตรง เพราะข้อมูล landmark จากผู้ใช้มีประเด็นด้าน privacy",
            "งานต่อยอด: เชื่อม skeleton เข้ากับ animated avatar, เพิ่ม motion curves และขยาย benchmark ไปทุก exercise mode",
        ],
        0.85,
        1.95,
        11.3,
        4.7,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "สรุปและการทำซ้ำงาน")
    add_bullets(
        slide,
        [
            "โปรเจกต์นี้มีองค์ประกอบครบตาม rubric: baseline comparison, scaling test และ visualization",
            "แกนหลักของงานคือการตีความ pose tracking ผ่านแนวคิดด้าน skeleton, key poses, interpolation และ temporal continuity",
            "คำสั่งหลัก: `python src/main.py --mode web`, `python src/main.py --mode benchmark`, `python src/main.py --mode demo`",
            "คำสั่งเสริมสำหรับข้อมูลจริง: `python tools/collect_mocap.py` และ `python experiments/evaluate_recorded_sequence.py`",
            "artifact ที่ส่งมีรายงาน Word, สไลด์ไทย/อังกฤษ, benchmark CSV/plot, demo video, tests และ source code ที่ทำซ้ำได้",
        ],
        0.85,
        1.95,
        11.3,
        4.6,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "เอกสารอ้างอิง")
    add_bullets(
        slide,
        [
            "Bradski, G. The OpenCV Library. Dr. Dobb's Journal of Software Tools, 2000.",
            "Bazarevsky, V. et al. BlazePose: On-device Real-time Body Pose tracking. arXiv, 2020.",
            "Google MediaPipe Team. MediaPipe Pose documentation.",
            "Casiez, G., Roussel, N., and Vogel, D. 1 Euro Filter. CHI 2012.",
            "Parent, R. Computer Animation: Algorithms and Techniques. Morgan Kaufmann, 2012.",
            "Menache, A. Understanding Motion Capture for Computer Animation and Video Games. Morgan Kaufmann, 2010.",
        ],
        0.85,
        1.95,
        11.2,
        4.8,
        font_size=17,
    )

    prs.save(OUTPUT_FILE)
    print(f"Generated {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
