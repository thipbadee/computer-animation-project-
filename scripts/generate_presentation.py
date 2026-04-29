from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from submission_metrics import (
    AUTHOR_NAME,
    COURSE_NAME,
    DEPARTMENT,
    GITHUB_USERNAME,
    INSTRUCTOR_NAME,
    KEYWORDS,
    PROJECT_TITLE,
    SEMESTER,
    benchmark_summary,
    group_rows_by_size,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_FILE = REPO_ROOT / "report" / "final_presentation.pptx"

TITLE_COLOR = RGBColor(20, 24, 32)
ACCENT_COLOR = RGBColor(0, 109, 119)
ACCENT_ALT = RGBColor(232, 248, 249)
TEXT_COLOR = RGBColor(50, 58, 70)
MUTED_COLOR = RGBColor(104, 116, 132)
WHITE = RGBColor(255, 255, 255)
LINE_COLOR = RGBColor(210, 219, 230)

WORKFLOW_IMAGE = REPO_ROOT / "report" / "assets" / "workflow.png"
PLOT_IMAGE = REPO_ROOT / "experiments" / "results" / "plots.png"
DEMO_STRIP_IMAGE = REPO_ROOT / "report" / "assets" / "demo_strip.png"


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.16), Inches(2.3), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT_COLOR
    rule.line.fill.background()

    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.28), Inches(12.0), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED_COLOR


def add_text(slide, text: str, left: float, top: float, width: float, height: float, font_size=18, bold=False, color=TEXT_COLOR):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR


def add_card(slide, title: str, value: str, left: float, top: float, width: float = 2.2, height: float = 1.1):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ALT
    shape.line.color.rgb = ACCENT_COLOR
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = ACCENT_COLOR
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = "Aptos Display"
    r2.font.size = Pt(18)
    r2.font.bold = True
    r2.font.color.rgb = TITLE_COLOR


def add_table(slide, data: list[list[str]], left: float, top: float, width: float, height: float, font_size=12):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    col_width = int((Inches(width) / cols))
    for idx in range(cols):
        table.columns[idx].width = col_width
    for r, row_data in enumerate(data):
        for c, value in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_ALT if r == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = TITLE_COLOR if r == 0 else TEXT_COLOR
                    run.font.bold = (r == 0)


def add_image(slide, image_path: Path, left: float, top: float, width: float):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def build_presentation():
    summary = benchmark_summary()
    grouped = group_rows_by_size()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, RGBColor(248, 250, 252))
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_COLOR
    banner.line.fill.background()
    add_text(slide, PROJECT_TITLE, 0.8, 1.05, 8.6, 0.85, font_size=28, bold=True, color=TITLE_COLOR)
    add_text(slide, COURSE_NAME, 0.8, 2.0, 2.8, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_text(slide, DEPARTMENT, 0.8, 2.35, 7.6, 0.45, font_size=14, color=MUTED_COLOR)
    add_text(slide, f"Author: {AUTHOR_NAME}", 0.8, 5.5, 3.0, 0.3, font_size=14)
    add_text(slide, f"GitHub: {GITHUB_USERNAME}", 0.8, 5.85, 3.2, 0.3, font_size=14)
    add_text(slide, f"Instructor: {INSTRUCTOR_NAME}", 0.8, 6.2, 3.4, 0.3, font_size=12, color=MUTED_COLOR)
    add_text(slide, f"Semester: {SEMESTER}", 4.0, 6.2, 3.6, 0.3, font_size=12, color=MUTED_COLOR)
    add_text(slide, "Keywords: " + ", ".join(KEYWORDS), 0.8, 6.55, 6.2, 0.3, font_size=12, color=MUTED_COLOR)
    add_card(slide, "100k Speedup", f"{summary['speedup_pct']:.2f}%", 8.8, 1.35)
    add_card(slide, "Memory Reduction", f"{summary['memory_reduction_pct']:.2f}%", 11.1, 1.35)
    add_card(slide, "Jitter Reduction", f"{summary['jitter_reduction_pct']:.2f}%", 9.95, 2.7)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Problem and Objectives", "Rubric coverage: problem significance, baseline/proposed scope, and engineering goals.")
    add_bullets(
        slide,
        [
            "Study real-time exercise motion tracking as a computer animation problem, not only a webcam classification task.",
            "Use pose landmarks as a simplified skeleton and interpret exercises as transitions between key poses.",
            "Compare a raw fixed-threshold baseline against a confidence-aware temporally filtered proposed method.",
            "Measure runtime, memory, repetition-count accuracy, and temporal stability under scaling.",
            "Deliver a reproducible artifact with live visualization, benchmark scripts, tests, and demo output.",
        ],
        0.8,
        1.9,
        6.1,
        4.8,
        font_size=18,
    )
    add_text(slide, "Why it fits Computer Animation", 7.3, 1.9, 2.8, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "Skeleton-like landmark representation",
            "Key poses and stage transitions",
            "Interpolation-derived progress gauge",
            "Markerless motion-capture interpretation",
        ],
        7.35,
        2.25,
        5.1,
        3.2,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Formal Problem Definition")
    add_bullets(
        slide,
        [
            "Input: time-ordered 2D joint landmarks with visibility from a webcam or recorded CSV.",
            "Output: repetition count, stage label, progress percentage, and overlay visualization.",
            "Constraints: interactive latency, occlusion tolerance, threshold robustness, and deterministic behavior.",
            "Challenge 1: noisy turning points cause count duplication or missed counts.",
            "Challenge 2: side switching creates unstable motion interpretation.",
            "Challenge 3: the live demo must remain understandable to the user in real time.",
        ],
        0.85,
        1.9,
        5.9,
        4.8,
        font_size=18,
    )
    add_text(slide, "Representative Motion Primitive", 7.1, 1.95, 3.3, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_text(
        slide,
        "The benchmark uses bicep curl tracking as a representative articulated motion primitive. The live app extends the same framework to Shoulder Press, Dumbbell Side Lateral Raise, Squat, Plank, and High Knees.",
        7.15,
        2.3,
        5.1,
        1.6,
        font_size=17,
    )
    add_text(slide, "Required Visualization", 7.1, 4.35, 2.6, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "Live Flask dashboard",
            "Side-by-side demo video",
            "Runtime and stability plots",
        ],
        7.15,
        4.7,
        4.4,
        1.8,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Baseline, Proposed, and SOTA Context")
    table_data = [
        ["Method", "Core Idea", "Strength", "Limitation"],
        ["Baseline", "Raw elbow angle + threshold count", "Very simple, deterministic", "Fragile under jitter and occlusion"],
        ["Proposed", "Confidence-aware side lock + EMA smoothing", "Lower jitter, faster benchmarked throughput", "Still handcrafted, not learned"],
        ["SOTA context", "Temporal neural / action models", "Richer sequence modeling", "Heavier compute and less explainable"],
    ]
    add_table(slide, table_data, 0.75, 1.9, 11.9, 2.4, font_size=11)
    add_bullets(
        slide,
        [
            "Both implemented methods keep the same input/output contract, so the comparison is fair.",
            "The proposed method targets engineering value: better stability while preserving real-time explainability.",
            "The project does not claim SOTA accuracy; it claims a reproducible improvement over the course-scale baseline.",
        ],
        0.9,
        4.75,
        11.0,
        1.8,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Proposed Algorithm and Workflow")
    add_image(slide, WORKFLOW_IMAGE, 6.9, 1.75, 5.7)
    add_bullets(
        slide,
        [
            "Compute left/right arm confidence from visibility and centeredness.",
            "Reuse the locked side when it remains reliable.",
            "Smooth the selected elbow angle with exponential filtering.",
            "Apply hysteresis thresholds to update stage and repetition count.",
            "Expose the same outputs as the baseline: count, stage, angle, progress.",
        ],
        0.85,
        1.95,
        5.5,
        3.6,
        font_size=18,
    )
    add_text(
        slide,
        "Complexity: O(1) time and O(1) memory per frame for both methods; the improvement comes from better state management, not larger asymptotic cost.",
        0.85,
        5.85,
        11.7,
        0.6,
        font_size=16,
        color=MUTED_COLOR,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Implementation Details")
    add_bullets(
        slide,
        [
            "Live demo stack: Flask + OpenCV + MediaPipe + NumPy.",
            "Benchmark stack: `src/algorithms/` for clean baseline/proposed comparison.",
            "Temporal state: counter, stage, active side, EMA filter value, side-lock timer.",
            "Edge cases handled: visibility drop, arm flapping, threshold chatter, camera read failure.",
            "Repository adds `tools/collect_mocap.py` and `experiments/evaluate_recorded_sequence.py` for real-data replay.",
        ],
        0.85,
        1.9,
        6.0,
        4.9,
        font_size=18,
    )
    add_text(slide, "Project Structure Highlights", 7.05, 1.9, 3.3, 0.3, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "`app.py`, `camera.py`, `exercises/`",
            "`src/algorithms/baseline.py`",
            "`src/algorithms/proposed.py`",
            "`experiments/benchmark.py`",
            "`visualization/animate_results.py`",
            "`tests/test_simulation.py`",
        ],
        7.1,
        2.25,
        5.0,
        3.6,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Experimental Setup")
    add_bullets(
        slide,
        [
            "Synthetic scaling benchmark sizes: 1k, 10k, and 100k frames.",
            "Synthetic generator injects angle noise and visibility dropouts to stress temporal robustness.",
            "Real-data path: webcam landmark capture to CSV, then deterministic replay for evaluation.",
            "Metrics: total runtime, ms/frame, peak traced memory, repetition count error, and side-switch jitter.",
            "Verification: `python -m unittest tests.test_simulation` passes on the current repository state.",
        ],
        0.85,
        1.95,
        6.1,
        4.6,
        font_size=18,
    )
    add_card(slide, "Synthetic Sizes", "1k / 10k / 100k", 7.4, 2.0, width=2.6)
    add_card(slide, "Count Error", "0 on all scales", 10.2, 2.0, width=2.3)
    add_card(slide, "Real-data Path", "Capture + Replay", 7.4, 3.45, width=2.6)
    add_card(slide, "Visualization", "Plots + Demo", 10.2, 3.45, width=2.3)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Quantitative Results")
    table_rows = [["Input", "Baseline", "Proposed", "Speedup", "Baseline Jitter", "Proposed Jitter"]]
    for size in sorted(grouped):
        baseline = grouped[size]["baseline"]
        proposed = grouped[size]["proposed"]
        speedup = ((baseline.ms_per_frame - proposed.ms_per_frame) / baseline.ms_per_frame) * 100.0
        table_rows.append(
            [
                f"{size:,}",
                f"{baseline.ms_per_frame:.6f}",
                f"{proposed.ms_per_frame:.6f}",
                f"{speedup:.2f}%",
                str(baseline.side_switches),
                str(proposed.side_switches),
            ]
        )
    add_table(slide, table_rows, 0.75, 1.9, 12.0, 2.5, font_size=11)
    add_card(slide, "100k ms/frame", f"{summary['baseline_ms_per_frame']:.6f} -> {summary['proposed_ms_per_frame']:.6f}", 0.95, 4.9, width=3.6)
    add_card(slide, "Peak Memory", f"{summary['baseline_peak_kb']:.3f} -> {summary['proposed_peak_kb']:.3f} KB", 4.95, 4.9, width=3.0)
    add_card(slide, "Side Switches", f"{int(summary['baseline_side_switches'])} -> {int(summary['proposed_side_switches'])}", 8.35, 4.9, width=3.1)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Plots and Visual Evaluation")
    add_image(slide, PLOT_IMAGE, 0.7, 1.75, 6.1)
    add_image(slide, DEMO_STRIP_IMAGE, 7.1, 1.9, 5.5)
    add_text(
        slide,
        "Left: quantitative scaling plots generated from `experiments/results/plots.png`. Right: representative side-by-side frames from the generated baseline/proposed demo video.",
        0.8,
        6.1,
        11.8,
        0.45,
        font_size=15,
        color=MUTED_COLOR,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Discussion, Limitations, and Future Work")
    add_bullets(
        slide,
        [
            "Strength: the proposed method improves temporal stability without changing the interface or requiring learned weights.",
            "Strength: the repository now includes benchmark scripts, tests, generated plots, demo video, and a mocap capture pipeline.",
            "Limitation: quantitative benchmark centers on one representative motion primitive rather than every live exercise mode.",
            "Limitation: the repository does not commit a human recording by default because webcam landmark data can be user-specific and privacy-sensitive.",
            "Future work: attach the tracked skeleton to an animated avatar, add richer motion curves, and benchmark all exercise modes.",
        ],
        0.85,
        1.95,
        11.3,
        4.7,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Conclusion and Reproducibility")
    add_bullets(
        slide,
        [
            "The project meets the course requirement of baseline comparison, scaling test, and visualization.",
            "It frames pose tracking as a computer animation artifact built from skeletons, key poses, interpolation, and temporal continuity.",
            "Main commands: `python src/main.py --mode web`, `python src/main.py --mode benchmark`, `python src/main.py --mode demo`.",
            "Additional pipeline: `python tools/collect_mocap.py` and `python experiments/evaluate_recorded_sequence.py`.",
            "Artifacts delivered: Word report, PowerPoint slides, benchmark CSV/plot, demo video, tests, and reproducible source code.",
        ],
        0.85,
        1.95,
        11.3,
        4.6,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "References")
    add_bullets(
        slide,
        [
            "Bradski, G. The OpenCV Library. Dr. Dobb's Journal of Software Tools, 2000.",
            "Bazarevsky, V. et al. BlazePose: On-device Real-time Body Pose tracking. arXiv, 2020.",
            "Google MediaPipe Team. MediaPipe Pose documentation.",
            "Casiez, G., Roussel, N., and Vogel, D. 1 Euro Filter. CHI 2012.",
            "Parent, R. Computer Animation: Algorithms and Techniques. Morgan Kaufmann, 2012.",
            "Menache, A. Understanding Motion Capture for Computer Animation and Video Games. Morgan Kaufmann, 2010.",
        ],
        0.85,
        1.95,
        11.2,
        4.8,
        font_size=17,
    )

    prs.save(OUTPUT_FILE)
    print(f"Generated {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
