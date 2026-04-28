from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


OUTPUT_FILE = "Computer_Animation_Project_Presentation.pptx"

TITLE_COLOR = RGBColor(31, 41, 55)
ACCENT_COLOR = RGBColor(37, 99, 235)
ACCENT_LIGHT = RGBColor(219, 234, 254)
TEXT_COLOR = RGBColor(55, 65, 81)
MUTED_COLOR = RGBColor(107, 114, 128)
PLACEHOLDER_FILL = RGBColor(243, 244, 246)
PLACEHOLDER_LINE = RGBColor(156, 163, 175)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(2.2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.32), Inches(11.3), Inches(0.55))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(12)
        r.font.color.rgb = MUTED_COLOR


def add_bullets(slide, items, left=0.8, top=1.8, width=5.4, height=4.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.bullet = True


def add_text_block(slide, text, left, top, width, height, font_size=18, bold=False, color=TEXT_COLOR):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_placeholder(slide, left, top, width, height, label, sublabel=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_FILL
    shape.line.color.rgb = PLACEHOLDER_LINE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = MUTED_COLOR

    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.name = "Aptos"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED_COLOR


def add_two_column_exercise_slide(slide, title, summary, cues):
    add_title(slide, title, "Leave the right panel for your own photo, demo frame, or posture illustration.")
    add_text_block(slide, "How to perform", 0.8, 1.8, 2.8, 0.4, font_size=19, bold=True, color=ACCENT_COLOR)
    add_text_block(slide, summary, 0.8, 2.2, 4.7, 1.25, font_size=18)
    add_text_block(slide, "Key points for presentation", 0.8, 3.55, 3.2, 0.35, font_size=18, bold=True, color=ACCENT_COLOR)
    add_bullets(slide, cues, left=0.95, top=3.9, width=4.7, height=2.8, font_size=17)
    add_placeholder(
        slide,
        6.2,
        1.9,
        6.5 - 1.0,
        4.5,
        "INSERT EXERCISE IMAGE HERE",
        "Suggested: photo of start pose and end pose, or one annotated screenshot.",
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(248, 250, 252)
    banner.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()
    add_text_block(
        slide,
        "Real-Time Exercise Tracking System",
        0.85,
        1.15,
        9.8,
        0.8,
        font_size=28,
        bold=True,
        color=TITLE_COLOR,
    )
    add_text_block(
        slide,
        "An Application of Computer Vision and Computer Animation",
        0.85,
        2.0,
        9.6,
        0.6,
        font_size=20,
        color=ACCENT_COLOR,
    )
    add_text_block(
        slide,
        "Computer Animation Course Presentation",
        0.85,
        2.6,
        5.5,
        0.4,
        font_size=16,
        color=MUTED_COLOR,
    )
    add_text_block(
        slide,
        "Prepared by: [Your Name] and [Teammate Name]",
        0.85,
        5.65,
        5.5,
        0.35,
        font_size=15,
        color=TEXT_COLOR,
    )
    add_text_block(
        slide,
        "Instructor: [Instructor Name]  |  Semester: [Semester / Academic Year]",
        0.85,
        6.05,
        6.8,
        0.35,
        font_size=13,
        color=MUTED_COLOR,
    )
    add_placeholder(
        slide,
        8.05,
        1.25,
        4.45,
        4.9,
        "OPTIONAL PROJECT IMAGE",
        "You may insert a dashboard screenshot, team logo, or a pose-tracking frame here.",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Project Overview", "This project combines computer vision with computer animation, with stronger emphasis on animation concepts.")
    add_bullets(
        slide,
        [
            "Build a real-time exercise tracking system using webcam input.",
            "Estimate body landmarks, detect movement stages, and count repetitions automatically.",
            "Interpret body motion as a sequence of poses and transitions rather than isolated frames.",
            "Use the project to demonstrate how animation theory can support motion analysis.",
        ],
    )
    add_placeholder(
        slide,
        7.2,
        1.8,
        5.2,
        4.8,
        "OPTIONAL SYSTEM SCREENSHOT",
        "Suggested: capture the web dashboard with pose landmarks and the progress overlay.",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Why This Project Fits Both Fields")
    add_text_block(slide, "Computer Animation", 0.85, 1.8, 2.8, 0.4, font_size=20, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "Represents the body as a simplified skeleton similar to a rig.",
            "Defines start poses, end poses, and in-between motion.",
            "Uses interpolation to estimate progress between key poses.",
            "Relates to motion capture, pose-to-pose animation, and temporal continuity.",
        ],
        left=0.95,
        top=2.2,
        width=5.4,
        height=3.7,
        font_size=17,
    )
    add_text_block(slide, "Computer Vision", 6.95, 1.8, 2.8, 0.4, font_size=20, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "Uses MediaPipe Pose to detect body landmarks from live video.",
            "Computes joint angles from landmark coordinates.",
            "Selects the clearer side of the body using landmark visibility.",
            "Provides the sensing layer for the animation-oriented interpretation.",
        ],
        left=7.05,
        top=2.2,
        width=5.2,
        height=3.7,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "System Workflow")
    workflow_steps = [
        ("1. Webcam Input", "Capture live video frames from the user."),
        ("2. Pose Estimation", "Detect body landmarks with MediaPipe Pose."),
        ("3. Motion Analysis", "Compute angles, stages, and progress for each exercise."),
        ("4. Animation Interpretation", "Treat the movement as key poses plus interpolation."),
        ("5. User Feedback", "Display reps, stage, angle, and gauge on the dashboard."),
    ]
    x_positions = [0.65, 3.15, 5.65, 8.15, 10.65]
    for idx, (heading, desc) in enumerate(workflow_steps):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(2.1), Inches(2.0), Inches(2.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_LIGHT
        shape.line.color.rgb = ACCENT_COLOR
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = heading
        r.font.name = "Aptos"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = ACCENT_COLOR
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = "Aptos"
        r2.font.size = Pt(12)
        r2.font.color.rgb = TEXT_COLOR
        if idx < len(workflow_steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x_positions[idx] + 2.05), Inches(2.78), Inches(0.35), Inches(0.6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.fill.background()

    add_text_block(
        slide,
        "Key presentation point: the sensing is done by computer vision, but the movement is interpreted through animation concepts.",
        0.9,
        5.25,
        11.4,
        0.6,
        font_size=17,
        color=MUTED_COLOR,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Computer Animation Concepts Applied", "This is the core emphasis of the project and should be stressed during the presentation.")
    add_bullets(
        slide,
        [
            "Skeletal representation: MediaPipe landmarks act like a simplified character rig.",
            "Key poses: each exercise has important states such as down, up, hold, or adjust.",
            "Interpolation: progress is estimated between the starting pose and ending pose.",
            "Temporal continuity: the system evaluates movement over time, not one frame only.",
            "Motion capture perspective: real human movement is converted into structured motion data.",
            "Readability of motion: the interface visualizes angle, progress, and stage clearly.",
        ],
        left=0.9,
        top=1.8,
        width=7.1,
        height=4.9,
        font_size=18,
    )
    add_placeholder(
        slide,
        8.55,
        1.95,
        3.9,
        4.45,
        "OPTIONAL DIAGRAM",
        "Suggested: insert a stick-figure skeleton, rig analogy, or pose transition illustration.",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Interpolation as the Main Animation Idea")
    add_text_block(slide, "Start Pose", 1.0, 2.1, 2.0, 0.3, font_size=22, bold=True, color=ACCENT_COLOR)
    add_text_block(slide, "Intermediate Motion", 4.9, 2.1, 3.0, 0.3, font_size=22, bold=True, color=ACCENT_COLOR)
    add_text_block(slide, "End Pose", 10.2, 2.1, 2.0, 0.3, font_size=22, bold=True, color=ACCENT_COLOR)
    add_placeholder(slide, 0.9, 2.6, 2.5, 2.7, "INSERT START POSE IMAGE", "Example: arm lowered or standing position")
    add_placeholder(slide, 5.0, 2.6, 3.2, 2.7, "INSERT IN-BETWEEN IMAGE", "Example: halfway through the movement")
    add_placeholder(slide, 10.0, 2.6, 2.5, 2.7, "INSERT END POSE IMAGE", "Example: arm raised or bottom position")
    add_text_block(
        slide,
        "The system uses joint-angle interpolation to estimate how far the user has moved from one key pose to another. "
        "This is why the progress gauge is not only a UI feature, but also an animation concept implemented in code.",
        0.95,
        5.7,
        11.2,
        0.75,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Role of Computer Vision", "Computer vision supports the project, but does not replace the animation-oriented interpretation.")
    add_bullets(
        slide,
        [
            "OpenCV handles video capture, frame processing, and screen overlay rendering.",
            "MediaPipe Pose estimates the positions of shoulders, elbows, wrists, hips, knees, and ankles.",
            "Visibility scores help choose the clearer body side when one side is occluded.",
            "Geometric angle calculation provides measurable input for stage detection.",
            "Threshold-based logic determines when a repetition has reached a meaningful pose state.",
        ],
        left=0.9,
        top=1.9,
        width=6.2,
        height=4.8,
        font_size=18,
    )
    add_placeholder(
        slide,
        7.8,
        2.0,
        4.5,
        4.4,
        "OPTIONAL POSE LANDMARK IMAGE",
        "Suggested: a frame showing MediaPipe's skeleton-style landmark output.",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Exercise Set in the Project")
    add_bullets(
        slide,
        [
            "Bicep Curl",
            "Shoulder Press",
            "Dumbbell Side Lateral Raise",
            "Squat",
            "Plank",
            "High Knees",
        ],
        left=1.0,
        top=2.0,
        width=4.0,
        height=3.8,
        font_size=22,
    )
    add_text_block(
        slide,
        "For presentation flow, each exercise can be explained through: start pose, end pose, intermediate motion, and what the system measures.",
        5.4,
        2.0,
        6.1,
        0.9,
        font_size=19,
    )
    add_placeholder(
        slide,
        6.0,
        3.05,
        5.8,
        3.1,
        "OPTIONAL COLLAGE OF EXERCISES",
        "You may insert one combined slide image showing all exercises or sample frames.",
    )

    exercise_slides = [
        (
            "Exercise 1: Bicep Curl",
            "Hold the dumbbell at your side, bend the elbow to raise it toward the shoulder, then lower it in a controlled manner.",
            [
                "Main measurement: elbow flexion and extension angle.",
                "Animation view: lowered-arm pose to flexed-arm pose.",
                "A full repetition occurs after a correct down-to-up transition.",
            ],
        ),
        (
            "Exercise 2: Shoulder Press",
            "Start with the weights near shoulder level, press the arms upward, then bring the weights down again with control.",
            [
                "Main measurement: arm extension during the press.",
                "Animation view: compressed pose to overhead extension pose.",
                "Useful to mention vertical motion and state transition.",
            ],
        ),
        (
            "Exercise 3: Dumbbell Side Lateral Raise",
            "Begin with the arms near the body, raise them outward to shoulder height, then lower them smoothly.",
            [
                "Main measurement: shoulder angle from the side of the torso to the raised arm.",
                "Animation view: clear start pose and end pose with easy interpolation.",
                "This is a strong example of animation-based motion progress.",
            ],
        ),
        (
            "Exercise 4: Squat",
            "Stand upright, lower the body by bending the hips and knees, then rise back to standing.",
            [
                "Main measurement: knee and leg-chain angle.",
                "Animation view: standing pose to lowered pose at the bottom.",
                "A good example of motion depth and posture control.",
            ],
        ),
        (
            "Exercise 5: Plank",
            "Hold the body in a straight supported position while keeping alignment from shoulders to ankles.",
            [
                "Main measurement: posture quality and holding time rather than repetition count.",
                "Animation view: maintaining one stable pose over time.",
                "Useful for explaining temporal continuity and pose stability.",
            ],
        ),
        (
            "Exercise 6: High Knees",
            "Run in place while alternately lifting the knees upward with a quick rhythm.",
            [
                "Main measurement: repeated alternating leg motion.",
                "Animation view: cyclical motion pattern rather than a single lift.",
                "Useful for explaining repeated motion timing.",
            ],
        ),
    ]

    for title, summary, cues in exercise_slides:
        slide = prs.slides.add_slide(blank)
        set_background(slide)
        add_two_column_exercise_slide(slide, title, summary, cues)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Implementation Summary")
    add_bullets(
        slide,
        [
            "Python project with modular exercise classes.",
            "Shared camera pipeline handles webcam capture and pose detection.",
            "Each exercise defines its own thresholds, stages, and counter logic.",
            "Dashboard displays the active exercise, current stage, and progress gauge.",
            "The interpolation-based percentage is both a practical metric and an animation concept.",
        ],
        left=0.9,
        top=1.9,
        width=6.4,
        height=4.8,
        font_size=18,
    )
    add_placeholder(
        slide,
        7.7,
        2.0,
        4.7,
        4.2,
        "OPTIONAL CODE / UI SNAPSHOT",
        "Suggested: insert dashboard UI, code structure, or exercise module screenshot.",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Results and Discussion")
    add_bullets(
        slide,
        [
            "The project shows that animation concepts can directly support interactive motion analysis.",
            "Users receive feedback that is easier to understand than simple classification alone.",
            "The strongest animation contribution is interpreting motion through skeletons, poses, and interpolation.",
            "The strongest vision contribution is automatic landmark acquisition from live camera input.",
            "Together, both fields create a practical system for exercise tracking and presentation of movement.",
        ],
        left=0.9,
        top=1.9,
        width=10.8,
        height=4.5,
        font_size=19,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "This project is a combination of computer vision and computer animation.",
            "However, its main academic value for this course lies in the animation-oriented interpretation of movement.",
            "The project applies skeletal representation, key poses, interpolation, motion continuity, and motion capture ideas.",
            "It demonstrates that computer animation can be used not only to create motion, but also to analyze real human motion.",
        ],
        left=0.9,
        top=1.95,
        width=10.7,
        height=4.7,
        font_size=20,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Future Improvements")
    add_bullets(
        slide,
        [
            "Attach the captured skeleton to a 2D or 3D animated avatar.",
            "Use smoother easing-style interpolation instead of only linear progress mapping.",
            "Store landmark sequences for playback and motion comparison.",
            "Add more exercises and more detailed pose checkpoints.",
            "Expand from analysis to real-time character animation demonstration.",
        ],
        left=0.9,
        top=1.95,
        width=10.8,
        height=4.7,
        font_size=20,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "References")
    add_bullets(
        slide,
        [
            "Bradski, G. The OpenCV Library.",
            "MediaPipe Pose documentation and pose landmark resources.",
            "Parent, R. Computer Animation: Algorithms and Techniques.",
            "Watt, A. and Watt, M. Advanced Animation and Rendering Techniques.",
            "Menache, A. Understanding Motion Capture for Computer Animation and Video Games.",
        ],
        left=0.9,
        top=1.95,
        width=10.9,
        height=4.8,
        font_size=18,
    )

    prs.save(OUTPUT_FILE)
    print(f"Generated {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
