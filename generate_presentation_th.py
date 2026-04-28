from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = "Computer_Animation_Project_Presentation_TH.pptx"

FONT_DISPLAY = "Leelawadee UI"
FONT_BODY = "Leelawadee UI"

TITLE_COLOR = RGBColor(31, 41, 55)
ACCENT_COLOR = RGBColor(37, 99, 235)
ACCENT_LIGHT = RGBColor(219, 234, 254)
TEXT_COLOR = RGBColor(55, 65, 81)
MUTED_COLOR = RGBColor(107, 114, 128)
PLACEHOLDER_FILL = RGBColor(243, 244, 246)
PLACEHOLDER_LINE = RGBColor(156, 163, 175)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.8), Inches(0.85))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_DISPLAY
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(2.2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(11.6), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = FONT_BODY
        r.font.size = Pt(12)
        r.font.color.rgb = MUTED_COLOR


def add_bullets(slide, items, left=0.8, top=1.8, width=5.6, height=4.8, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(7)
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.bullet = True


def add_text_block(slide, text, left, top, width, height, font_size=18, bold=False, color=TEXT_COLOR, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_placeholder(slide, left, top, width, height, label, sublabel=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_FILL
    shape.line.color.rgb = PLACEHOLDER_LINE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT_BODY
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = MUTED_COLOR

    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sublabel
        r2.font.name = FONT_BODY
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED_COLOR


def add_two_column_exercise_slide(slide, title, summary, cues):
    add_title(slide, title, "เว้นพื้นที่ด้านขวาไว้สำหรับใส่ภาพท่าออกกำลังกายจริง ภาพเริ่มต้น-สิ้นสุด หรือภาพอธิบายท่าทาง")
    add_text_block(slide, "วิธีเล่นโดยสรุป", 0.8, 1.8, 2.5, 0.35, font_size=18, bold=True, color=ACCENT_COLOR)
    add_text_block(slide, summary, 0.8, 2.15, 4.9, 1.35, font_size=17)
    add_text_block(slide, "ประเด็นที่ควรพูดตอนพรีเซนต์", 0.8, 3.6, 3.8, 0.35, font_size=17, bold=True, color=ACCENT_COLOR)
    add_bullets(slide, cues, left=0.95, top=3.95, width=4.9, height=2.6, font_size=16)
    add_placeholder(
        slide,
        6.15,
        1.95,
        5.6,
        4.45,
        "แทรกรูปท่าออกกำลังกายที่นี่",
        "แนะนำ: รูปท่าเริ่มต้น/ท่าสิ้นสุด หรือภาพ screenshot พร้อม landmark",
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    bg.line.fill.background()
    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.28))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_COLOR
    top_bar.line.fill.background()
    add_text_block(slide, "ระบบติดตามท่าออกกำลังกายแบบเรียลไทม์", 0.85, 1.1, 8.7, 0.7, font_size=28, bold=True, color=TITLE_COLOR)
    add_text_block(slide, "การประยุกต์ใช้ Computer Vision และ Computer Animation", 0.85, 1.95, 8.2, 0.55, font_size=20, color=ACCENT_COLOR)
    add_text_block(slide, "สไลด์นำเสนอรายวิชา Computer Animation", 0.85, 2.55, 5.6, 0.35, font_size=16, color=MUTED_COLOR)
    add_text_block(slide, "ผู้จัดทำ: [ชื่อผู้จัดทำ] และ [ชื่อเพื่อนร่วมทีม]", 0.85, 5.65, 5.8, 0.32, font_size=14, color=TEXT_COLOR)
    add_text_block(slide, "อาจารย์ผู้สอน: [ชื่ออาจารย์] | ภาคการศึกษา: [ภาคการศึกษา/ปีการศึกษา]", 0.85, 6.0, 7.2, 0.32, font_size=12, color=MUTED_COLOR)
    add_placeholder(
        slide,
        8.05,
        1.2,
        4.45,
        4.95,
        "ใส่ภาพโปรเจกต์ได้ที่นี่",
        "เช่น ภาพหน้า dashboard, ภาพ landmark skeleton หรือภาพรวมของระบบ",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "ภาพรวมของโปรเจกต์", "โปรเจกต์นี้เป็นการผสานกันของ Computer Vision และ Computer Animation แต่ให้น้ำหนักกับ Computer Animation มากกว่า")
    add_bullets(
        slide,
        [
            "พัฒนาระบบติดตามท่าออกกำลังกายแบบเรียลไทม์จากภาพกล้องเว็บแคม",
            "ตรวจจับตำแหน่งข้อต่อของร่างกาย วิเคราะห์สถานะการเคลื่อนไหว และนับจำนวนครั้งอัตโนมัติ",
            "ตีความการเคลื่อนไหวเป็นลำดับของ pose และการเปลี่ยนผ่าน ไม่ใช่มองเป็นเพียงภาพนิ่งรายเฟรม",
            "ใช้โปรเจกต์นี้เพื่อแสดงให้เห็นว่าแนวคิดทาง animation สามารถประยุกต์สู่การวิเคราะห์การเคลื่อนไหวได้",
        ],
        left=0.9,
        top=1.9,
        width=6.0,
        height=4.6,
        font_size=18,
    )
    add_placeholder(
        slide,
        7.4,
        1.9,
        5.0,
        4.7,
        "ใส่ภาพหน้าระบบที่นี่",
        "แนะนำ: ภาพหน้าเว็บพร้อม video stream, skeleton landmark, และ progress gauge",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "เหตุผลที่โปรเจกต์นี้เกี่ยวข้องกับทั้งสองศาสตร์")
    add_text_block(slide, "Computer Animation", 0.85, 1.8, 3.0, 0.4, font_size=20, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "แทนร่างกายมนุษย์ด้วยโครงกระดูกแบบย่อส่วน คล้าย rig ของตัวละคร",
            "กำหนดท่าเริ่มต้น ท่าสิ้นสุด และช่วงการเคลื่อนไหวระหว่างกลาง",
            "ใช้ interpolation เพื่อประมาณความคืบหน้าระหว่าง key poses",
            "เชื่อมโยงกับ motion capture, pose-to-pose animation และ temporal continuity",
        ],
        left=0.95,
        top=2.2,
        width=5.35,
        height=3.9,
        font_size=16,
    )
    add_text_block(slide, "Computer Vision", 6.95, 1.8, 2.8, 0.4, font_size=20, bold=True, color=ACCENT_COLOR)
    add_bullets(
        slide,
        [
            "ใช้ MediaPipe Pose ตรวจจับจุด landmark ของร่างกายจากวิดีโอ",
            "คำนวณมุมของข้อต่อจากพิกัดที่ตรวจจับได้",
            "ใช้ visibility ช่วยเลือกด้านของร่างกายที่เห็นชัดกว่า",
            "ทำหน้าที่เป็น sensing layer ให้กับการตีความเชิง animation",
        ],
        left=7.05,
        top=2.2,
        width=5.1,
        height=3.9,
        font_size=16,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "ลำดับการทำงานของระบบ")
    workflow_steps = [
        ("1. รับภาพจากกล้อง", "ดึงภาพวิดีโอจากผู้ใช้แบบเรียลไทม์"),
        ("2. ตรวจจับท่าทาง", "หา landmark ของร่างกายด้วย MediaPipe Pose"),
        ("3. วิเคราะห์การเคลื่อนไหว", "คำนวณมุม, stage และ progress ของแต่ละท่า"),
        ("4. ตีความเชิง Animation", "มองการเคลื่อนไหวเป็น key poses และ in-between motion"),
        ("5. แสดงผล", "โชว์ reps, stage, angle และ gauge บนหน้าระบบ"),
    ]
    x_positions = [0.62, 3.14, 5.66, 8.18, 10.7]
    for idx, (heading, desc) in enumerate(workflow_steps):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(2.05), Inches(2.0), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_LIGHT
        shape.line.color.rgb = ACCENT_COLOR
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = heading
        r.font.name = FONT_BODY
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = ACCENT_COLOR
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT_BODY
        r2.font.size = Pt(11)
        r2.font.color.rgb = TEXT_COLOR
        if idx < len(workflow_steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x_positions[idx] + 2.05), Inches(2.88), Inches(0.34), Inches(0.56))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.fill.background()
    add_text_block(
        slide,
        "ประเด็นสำคัญที่ควรพูด: ส่วนที่ทำหน้าที่รับรู้ข้อมูลคือ computer vision แต่การตีความ movement progression เป็นจุดเด่นของ computer animation",
        0.9,
        5.25,
        11.4,
        0.65,
        font_size=16,
        color=MUTED_COLOR,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "แนวคิดด้าน Computer Animation ที่ใช้ในโปรเจกต์", "สไลด์นี้เป็นแกนหลักของงานพรีเซนต์ ควรเน้นเป็นพิเศษ")
    add_bullets(
        slide,
        [
            "Skeletal representation: landmark ของ MediaPipe ทำหน้าที่คล้าย rig ของตัวละคร",
            "Key poses: แต่ละท่ามีสถานะสำคัญ เช่น down, up, hold หรือ adjust",
            "Interpolation: ระบบประมาณความคืบหน้าจากท่าเริ่มต้นไปท่าสิ้นสุด",
            "Temporal continuity: พิจารณาการเคลื่อนไหวต่อเนื่องตามเวลา ไม่ใช่เฟรมเดียว",
            "Motion capture perspective: นำการเคลื่อนไหวของคนจริงมาแทนเป็นข้อมูลเชิงโครงกระดูก",
            "Readability of motion: แสดงค่า stage, angle และ progress ให้เข้าใจง่าย",
        ],
        left=0.9,
        top=1.85,
        width=7.15,
        height=4.9,
        font_size=17,
    )
    add_placeholder(
        slide,
        8.5,
        1.95,
        3.95,
        4.45,
        "ใส่ภาพประกอบแนวคิดได้ที่นี่",
        "เช่น ภาพ stick figure, rig analogy, หรือภาพการเปลี่ยน pose",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "Interpolation: แนวคิด Animation ที่เด่นที่สุด")
    add_text_block(slide, "ท่าเริ่มต้น", 1.05, 2.05, 2.0, 0.3, font_size=21, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
    add_text_block(slide, "ท่าระหว่างทาง", 4.95, 2.05, 2.8, 0.3, font_size=21, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
    add_text_block(slide, "ท่าสิ้นสุด", 10.05, 2.05, 2.0, 0.3, font_size=21, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
    add_placeholder(slide, 0.95, 2.55, 2.45, 2.7, "ใส่รูปท่าเริ่มต้น", "เช่น ท่ายืนตรง หรือแขนอยู่ด้านล่าง")
    add_placeholder(slide, 4.85, 2.55, 3.1, 2.7, "ใส่รูปท่าระหว่างทาง", "เช่น ภาพช่วงครึ่งทางของการเคลื่อนไหว")
    add_placeholder(slide, 9.9, 2.55, 2.45, 2.7, "ใส่รูปท่าสิ้นสุด", "เช่น ท่าแขนยกขึ้น หรือท่าย่อลงสุด")
    add_text_block(
        slide,
        "ระบบใช้การคำนวณมุมของข้อต่อและนำไป map เป็นเปอร์เซ็นต์ความคืบหน้า เพื่อบอกว่าผู้ใช้อยู่ในช่วงใดระหว่าง key pose แรกและ key pose สุดท้าย ซึ่งสะท้อนแนวคิด interpolation โดยตรง",
        0.95,
        5.65,
        11.3,
        0.72,
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "บทบาทของ Computer Vision", "เป็นส่วนสนับสนุนระบบ แต่การอธิบายเชิงวิชาควรเชื่อมกลับไปยัง animation")
    add_bullets(
        slide,
        [
            "OpenCV ใช้สำหรับรับภาพจากกล้อง ประมวลผลเฟรม และแสดง overlay บนหน้าจอ",
            "MediaPipe Pose ใช้หา landmark ของไหล่ ศอก ข้อมือ สะโพก เข่า และข้อเท้า",
            "ค่า visibility ช่วยให้เลือกร่างกายด้านที่เห็นชัดกว่าเมื่อตรวจจับท่า",
            "การคำนวณมุมจาก 3 จุด ทำให้ระบบวัดสถานะของการเคลื่อนไหวได้",
            "threshold-based logic ใช้ตัดสินว่าถึงท่าสำคัญสำหรับการนับ rep หรือยัง",
        ],
        left=0.9,
        top=1.9,
        width=6.35,
        height=4.8,
        font_size=17,
    )
    add_placeholder(
        slide,
        7.85,
        2.0,
        4.5,
        4.35,
        "ใส่ภาพ skeleton landmark ได้ที่นี่",
        "แนะนำ: ภาพจากระบบที่เห็นโครงก้างปลาของ MediaPipe ชัดเจน",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "ชุดท่าออกกำลังกายในโปรเจกต์")
    add_bullets(
        slide,
        [
            "Bicep Curl",
            "Shoulder Press",
            "Dumbbell Side Lateral Raise",
            "Squat",
            "Plank",
            "High Knees",
        ],
        left=1.0,
        top=2.0,
        width=4.0,
        height=3.8,
        font_size=21,
    )
    add_text_block(
        slide,
        "เวลาพรีเซนต์แต่ละท่า ควรอธิบาย 4 ส่วน: ท่าเริ่มต้น, ท่าสิ้นสุด, ช่วงการเปลี่ยนผ่าน, และสิ่งที่ระบบใช้วัดการเคลื่อนไหว",
        5.45,
        2.0,
        6.0,
        0.95,
        font_size=18,
    )
    add_placeholder(
        slide,
        6.0,
        3.05,
        5.8,
        3.1,
        "ใส่ภาพรวมของทุกท่าได้ที่นี่",
        "เช่น collage ของท่าออกกำลังกายทั้งหมด หรือภาพ sample หลายเฟรมรวมกัน",
    )

    exercise_slides = [
        (
            "ท่าที่ 1: Bicep Curl",
            "ถือดัมเบลไว้ข้างลำตัว จากนั้นงอศอกยกดัมเบลขึ้นเข้าหาไหล่ แล้วค่อยลดลงอย่างควบคุม",
            [
                "ระบบวัดมุมการงอและเหยียดของข้อศอก",
                "มองเชิง animation เป็นการเปลี่ยนจากท่าแขนลงไปสู่ท่าแขนงอขึ้น",
                "นับ 1 ครั้งเมื่อเกิดลำดับการเคลื่อนไหวจาก down ไป up อย่างถูกต้อง",
            ],
        ),
        (
            "ท่าที่ 2: Shoulder Press",
            "เริ่มจากถือดัมเบลบริเวณระดับไหล่ แล้วดันแขนขึ้นด้านบนก่อนลดกลับลงมาอย่างควบคุม",
            [
                "ระบบวัดการเหยียดแขนในจังหวะดันขึ้น",
                "มองเชิง animation เป็นการเปลี่ยนจากท่าหดตัวไปสู่ท่าเหยียดขึ้นเหนือศีรษะ",
                "เหมาะกับการอธิบายการเคลื่อนที่แนวดิ่งและ state transition",
            ],
        ),
        (
            "ท่าที่ 3: Dumbbell Side Lateral Raise",
            "เริ่มจากวางแขนไว้ข้างลำตัว แล้วยกแขนออกด้านข้างจนใกล้ระดับไหล่ ก่อนลดลงอย่างนุ่มนวล",
            [
                "ระบบวัดมุมหัวไหล่ระหว่างลำตัวกับแขนที่ยกออกด้านข้าง",
                "เป็นตัวอย่างที่ชัดมากของท่าเริ่มต้น ท่าสิ้นสุด และ interpolation ระหว่างทาง",
                "ท่านี้เหมาะที่สุดสำหรับเชื่อมกับแนวคิด animation ในรายงาน",
            ],
        ),
        (
            "ท่าที่ 4: Squat",
            "ยืนตรงก่อน จากนั้นย่อตัวลงโดยงอสะโพกและเข่า แล้วดันตัวกลับขึ้นมายืนตรง",
            [
                "ระบบวัดมุมของขาและเข่าเพื่อดูความลึกของการย่อ",
                "มองเชิง animation เป็นการเปลี่ยนจากท่ายืนไปสู่ท่าย่อลงสุด",
                "เหมาะกับการอธิบายเรื่อง motion depth และ posture control",
            ],
        ),
        (
            "ท่าที่ 5: Plank",
            "ค้างลำตัวให้ตรงจากไหล่ไปถึงข้อเท้า โดยพยายามไม่ยกสะโพกสูงหรือปล่อยลำตัวตก",
            [
                "ระบบเน้นวัดคุณภาพของ posture และระยะเวลาที่ค้างท่า",
                "ในเชิง animation คือการรักษา pose เดิมให้คงที่ตลอดช่วงเวลา",
                "ใช้เชื่อมกับแนวคิด temporal continuity และ pose stability ได้ดี",
            ],
        ),
        (
            "ท่าที่ 6: High Knees",
            "วิ่งอยู่กับที่พร้อมยกเข่าสลับซ้ายขวาอย่างต่อเนื่องและมีจังหวะ",
            [
                "ระบบวัดรูปแบบการเคลื่อนไหวสลับของขาทั้งสองข้าง",
                "ในเชิง animation เป็นการเคลื่อนไหวแบบ cyclical motion",
                "เหมาะกับการอธิบายจังหวะและการเคลื่อนไหวซ้ำ ๆ ตามเวลา",
            ],
        ),
    ]

    for title, summary, cues in exercise_slides:
        slide = prs.slides.add_slide(blank)
        set_background(slide)
        add_two_column_exercise_slide(slide, title, summary, cues)

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "สรุปการพัฒนาระบบ")
    add_bullets(
        slide,
        [
            "โปรเจกต์เขียนด้วย Python และแยก exercise logic เป็นคลาสของแต่ละท่า",
            "มี camera pipeline กลางสำหรับรับภาพและตรวจจับ pose",
            "แต่ละท่ากำหนด threshold, stage และ counter logic ของตนเอง",
            "หน้า dashboard ใช้แสดง active exercise, stage ปัจจุบัน และ progress gauge",
            "เปอร์เซ็นต์ความคืบหน้าไม่ใช่แค่ UI แต่เป็นการใช้แนวคิด interpolation จริงในระบบ",
        ],
        left=0.9,
        top=1.95,
        width=6.45,
        height=4.8,
        font_size=17,
    )
    add_placeholder(
        slide,
        7.75,
        2.0,
        4.6,
        4.25,
        "ใส่ภาพโค้ดหรือ UI ได้ที่นี่",
        "แนะนำ: โครงสร้างโค้ด exercise, ภาพ overlay, หรือหน้าเว็บของระบบ",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "ผลลัพธ์และการอภิปราย")
    add_bullets(
        slide,
        [
            "โปรเจกต์แสดงให้เห็นว่าแนวคิดทาง animation สามารถนำมาช่วยวิเคราะห์การเคลื่อนไหวได้จริง",
            "ผู้ใช้ได้รับ feedback ที่เข้าใจง่ายกว่าการบอกเพียงว่า 'ถูก' หรือ 'ผิด'",
            "จุดเด่นด้าน animation คือการมอง motion เป็น skeleton, key pose และ in-between estimation",
            "จุดเด่นด้าน computer vision คือการได้ landmark จากกล้องแบบอัตโนมัติ",
            "การรวมกันของทั้งสองศาสตร์ทำให้เกิดระบบติดตามการออกกำลังกายที่ใช้งานได้จริง",
        ],
        left=0.9,
        top=1.95,
        width=10.9,
        height=4.6,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "บทสรุป")
    add_bullets(
        slide,
        [
            "โปรเจกต์นี้เป็นการบูรณาการระหว่าง Computer Vision และ Computer Animation",
            "อย่างไรก็ตาม คุณค่าหลักของงานในบริบทวิชาอยู่ที่การตีความการเคลื่อนไหวแบบ animation-oriented",
            "ระบบใช้แนวคิด skeletal representation, key poses, interpolation, temporal continuity และ motion capture",
            "จึงแสดงให้เห็นว่า computer animation ไม่ได้ใช้เพียงเพื่อสร้างภาพเคลื่อนไหว แต่ยังใช้วิเคราะห์การเคลื่อนไหวจริงของมนุษย์ได้ด้วย",
        ],
        left=0.9,
        top=1.95,
        width=10.8,
        height=4.8,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "แนวทางพัฒนาต่อ")
    add_bullets(
        slide,
        [
            "เชื่อม skeleton ที่ตรวจจับได้เข้ากับตัวละคร 2D หรือ 3D",
            "พัฒนา interpolation ให้มีลักษณะ easing มากขึ้น แทน linear progress แบบง่าย",
            "บันทึกลำดับ landmark เพื่อเล่นย้อนหลังและเปรียบเทียบ motion ได้",
            "เพิ่มจำนวนท่าออกกำลังกายและเพิ่ม pose checkpoint ให้ละเอียดขึ้น",
            "ขยายจากการวิเคราะห์ movement ไปสู่การแสดง character animation แบบเรียลไทม์",
        ],
        left=0.9,
        top=1.95,
        width=10.8,
        height=4.8,
        font_size=18,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "เอกสารอ้างอิง")
    add_bullets(
        slide,
        [
            "Bradski, G. The OpenCV Library.",
            "MediaPipe Pose documentation และแหล่งข้อมูลเกี่ยวกับ pose landmark estimation.",
            "Parent, R. Computer Animation: Algorithms and Techniques.",
            "Watt, A. และ Watt, M. Advanced Animation and Rendering Techniques.",
            "Menache, A. Understanding Motion Capture for Computer Animation and Video Games.",
        ],
        left=0.9,
        top=1.95,
        width=10.95,
        height=4.7,
        font_size=17,
    )

    prs.save(OUTPUT_FILE)
    print(f"Generated {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
